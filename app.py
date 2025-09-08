"""
FAISS Vector Search Educational App
===================================

This Streamlit app demonstrates how FAISS (Facebook AI Similarity Search) works
for vector-based document retrieval. It provides an interactive interface to:
1. Add documents and see their embeddings
2. Search for similar documents
3. Visualize how FAISS finds nearest neighbors

Author: AI Assistant
"""

import streamlit as st
import numpy as np
import faiss
import matplotlib.pyplot as plt
import plotly.express as px
import plotly.graph_objects as go
from sentence_transformers import SentenceTransformer
from sklearn.decomposition import PCA
from sklearn.manifold import TSNE
import pandas as pd
from typing import List, Tuple, Optional
import time
import io

# File processing imports with better error handling
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PyPDF2 = None
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    Document = None
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None
    PPTX_AVAILABLE = False

# Configure page
st.set_page_config(
    page_title="FAISS Vector Search Demo",
    page_icon="🔍",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better styling
st.markdown("""
<style>
    .section-container {
        background-color: #f8f9fa;
        padding: 1.5rem;
        border-radius: 10px;
        margin-bottom: 1rem;
        border: 1px solid #e9ecef;
    }
    
    .input-container {
        background-color: white;
        padding: 1rem;
        border-radius: 8px;
        border: 1px solid #dee2e6;
        margin-bottom: 1rem;
    }
    
    .document-container {
        background-color: #ffffff;
        padding: 1rem;
        border-radius: 8px;
        border-left: 4px solid #007bff;
        margin-bottom: 0.5rem;
        box-shadow: 0 1px 3px rgba(0,0,0,0.1);
    }
    
    .search-result {
        background-color: #f8f9fa;
        padding: 1rem;
        border-radius: 8px;
        border-left: 4px solid #28a745;
        margin-bottom: 1rem;
        box-shadow: 0 1px 3px rgba(0,0,0,0.1);
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state
if 'documents' not in st.session_state:
    st.session_state.documents = []
if 'embeddings' not in st.session_state:
    st.session_state.embeddings = []
if 'faiss_index' not in st.session_state:
    st.session_state.faiss_index = None
if 'model' not in st.session_state:
    st.session_state.model = None
if 'search_results' not in st.session_state:
    st.session_state.search_results = None

@st.cache_resource
def load_model():
    """Load the sentence transformer model"""
    return SentenceTransformer('all-MiniLM-L6-v2')

def embed_text(text: str, model: SentenceTransformer) -> np.ndarray:
    """
    Convert text to embedding vector using sentence transformers
    
    Args:
        text: Input text to embed
        model: Pre-loaded sentence transformer model
        
    Returns:
        Embedding vector as numpy array
    """
    embedding = model.encode([text])
    return embedding[0]

def extract_text_from_pdf(file) -> str:
    """
    Extract text from uploaded PDF file
    
    Args:
        file: Streamlit uploaded file object
        
    Returns:
        Extracted text as string
    """
    if not PDF_AVAILABLE:
        return "PDF processing not available. Please install PyPDF2 or contact support."
    
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file.read()))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        return f"Error reading PDF: {str(e)}"

def extract_text_from_docx(file) -> str:
    """
    Extract text from uploaded Word document
    
    Args:
        file: Streamlit uploaded file object
        
    Returns:
        Extracted text as string
    """
    if not DOCX_AVAILABLE:
        return "Word document processing not available. Please install python-docx or contact support."
    
    try:
        doc = Document(io.BytesIO(file.read()))
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        return f"Error reading Word document: {str(e)}"

def extract_text_from_pptx(file) -> str:
    """
    Extract text from uploaded PowerPoint presentation
    
    Args:
        file: Streamlit uploaded file object
        
    Returns:
        Extracted text as string
    """
    if not PPTX_AVAILABLE:
        return "PowerPoint processing not available. Please install python-pptx or contact support."
    
    try:
        prs = Presentation(io.BytesIO(file.read()))
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"Slide {slide_num}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
            text += "\n"
        return text.strip()
    except Exception as e:
        return f"Error reading PowerPoint: {str(e)}"

def process_uploaded_file(file) -> Tuple[str, str]:
    """
    Process uploaded file and extract text based on file type
    
    Args:
        file: Streamlit uploaded file object
        
    Returns:
        Tuple of (extracted_text, file_info)
    """
    file_extension = file.name.lower().split('.')[-1]
    file_info = f"📄 {file.name} ({file.size:,} bytes)"
    
    if file_extension == 'pdf':
        text = extract_text_from_pdf(file)
    elif file_extension in ['docx', 'doc']:
        text = extract_text_from_docx(file)
    elif file_extension in ['pptx', 'ppt']:
        text = extract_text_from_pptx(file)
    else:
        text = f"Unsupported file type: .{file_extension}"
    
    return text, file_info

def create_faiss_index(embeddings: List[np.ndarray]) -> faiss.IndexFlatIP:
    """
    Create a FAISS index for inner product (cosine similarity) search
    
    Args:
        embeddings: List of embedding vectors
        
    Returns:
        FAISS index object
    """
    if not embeddings:
        return None
    
    # Stack embeddings into a matrix
    embedding_matrix = np.vstack(embeddings).astype('float32')
    
    # Normalize vectors for cosine similarity
    faiss.normalize_L2(embedding_matrix)
    
    # Create FAISS index (Inner Product for normalized vectors = cosine similarity)
    dimension = embedding_matrix.shape[1]
    index = faiss.IndexFlatIP(dimension)
    
    # Add vectors to index
    index.add(embedding_matrix)
    
    return index

def search_faiss(query_embedding: np.ndarray, index: faiss.IndexFlatIP, k: int = 5) -> Tuple[np.ndarray, np.ndarray]:
    """
    Search for similar vectors in FAISS index
    
    Args:
        query_embedding: Query vector
        index: FAISS index
        k: Number of nearest neighbors to return
        
    Returns:
        Tuple of (similarity_scores, indices)
    """
    if index is None:
        return np.array([]), np.array([])
    
    # Normalize query vector
    query_norm = query_embedding.reshape(1, -1).astype('float32')
    faiss.normalize_L2(query_norm)
    
    # Search
    scores, indices = index.search(query_norm, min(k, index.ntotal))
    
    return scores[0], indices[0]

def reduce_dimensions(embeddings: List[np.ndarray], method: str = 'pca') -> np.ndarray:
    """
    Reduce embedding dimensions to 2D for visualization
    
    Args:
        embeddings: List of embedding vectors
        method: 'pca' or 'tsne'
        
    Returns:
        2D coordinates as numpy array
    """
    if len(embeddings) < 2:
        return np.array([[0, 0]] * len(embeddings))
    
    embedding_matrix = np.vstack(embeddings)
    
    if method == 'pca':
        reducer = PCA(n_components=2, random_state=42)
    else:  # tsne
        reducer = TSNE(n_components=2, random_state=42, perplexity=min(30, len(embeddings)-1))
    
    return reducer.fit_transform(embedding_matrix)

def plot_embeddings_interactive(embeddings_2d: np.ndarray, documents: List[str], 
                               query_2d: Optional[np.ndarray] = None, 
                               similar_indices: Optional[List[int]] = None,
                               similarity_scores: Optional[np.ndarray] = None) -> go.Figure:
    """
    Create interactive 2D scatter plot of embeddings using Plotly
    
    Args:
        embeddings_2d: 2D coordinates of document embeddings
        documents: List of document texts
        query_2d: 2D coordinates of query (optional)
        similar_indices: Indices of similar documents to highlight
        
    Returns:
        Plotly figure object
    """
    if len(embeddings_2d) == 0:
        fig = go.Figure()
        fig.update_layout(title="No documents added yet")
        return fig
    
    # Calculate sizes based on similarity scores if available
    sizes = []
    colors = []
    
    for i in range(len(documents)):
        if i in (similar_indices or []) and similarity_scores is not None:
            # Find the similarity score for this document
            idx_pos = list(similar_indices).index(i) if i in similar_indices else -1
            if idx_pos >= 0 and idx_pos < len(similarity_scores):
                # Scale size based on similarity (0.3 to 1.0 similarity -> 15 to 25 size)
                score = similarity_scores[idx_pos]
                size = 15 + (score * 10)  # Size between 15-25 based on similarity
            else:
                size = 15
            colors.append('Similar')
        else:
            # Non-similar documents are smaller and blue
            size = 10
            colors.append('Document')
        sizes.append(size)
    
    # Create DataFrame for easier plotting
    df = pd.DataFrame({
        'x': embeddings_2d[:, 0],
        'y': embeddings_2d[:, 1],
        'text': [f"Doc {i+1}: {doc[:50]}..." if len(doc) > 50 else f"Doc {i+1}: {doc}" 
                for i, doc in enumerate(documents)],
        'color': colors,
        'size': sizes
    })
    
    fig = px.scatter(df, x='x', y='y', color='color', size='size',
                    hover_data=['text'],
                    title="📊 Document Embeddings in 2D Vector Space",
                    color_discrete_map={'Document': '#4A90E2', 'Similar': '#E74C3C'})
    
    # Add query point if provided
    if query_2d is not None:
        fig.add_trace(go.Scatter(
            x=[query_2d[0]], y=[query_2d[1]],
            mode='markers',
            marker=dict(color='#2ECC71', size=20, symbol='star', line=dict(width=2, color='white')),
            name='🔍 Your Query',
            hovertext=['Your Search Query'],
            showlegend=True
        ))
    
    fig.update_layout(
        width=800, height=600,
        showlegend=True,
        xaxis_title="Dimension 1",
        yaxis_title="Dimension 2"
    )
    
    return fig

def display_educational_panel():
    """Display educational information about FAISS"""
    st.sidebar.header("🎓 How FAISS Works")
    
    with st.sidebar.expander("What is FAISS?", expanded=True):
        st.write("""
        **FAISS** (Facebook AI Similarity Search) is a library for efficient similarity search 
        and clustering of dense vectors.
        
        **Key Concepts:**
        - 📝 **Documents** → **Vectors**: Text is converted to high-dimensional vectors (embeddings)
        - 📁 **File Support**: Upload PDF, Word, PowerPoint files for automatic text extraction
        - 📏 **Similarity**: Measures how "close" vectors are in the embedding space
        - 🔍 **Search**: Finds the most similar vectors to your query
        """)
    
    with st.sidebar.expander("The Search Process"):
        st.write("""
        1. **Embedding**: Convert your query text to a vector
        2. **Normalize**: Prepare vector for cosine similarity
        3. **Search**: FAISS finds nearest neighbors efficiently
        4. **Rank**: Results sorted by similarity score
        
        **Similarity Score**: Higher = more similar (max = 1.0)
        """)
    
    with st.sidebar.expander("📊 Visualization Guide"):
        st.write("""
        **Understanding the 2D Plot:**
        - 🔵 **Blue dots**: All your documents in vector space
        - ⭐ **Green star**: Your search query position
        - 🔴 **Red dots**: Documents similar to your query
        - **Dot size**: Larger red dots = higher similarity scores
        
        **Technical Note**: We use PCA/t-SNE to reduce 384D embeddings to 2D.
        Some spatial relationships may change, but semantic clusters remain visible!
        """)

def main():
    """Main application function"""
    
    # Load model
    if st.session_state.model is None:
        with st.spinner("Loading sentence transformer model..."):
            st.session_state.model = load_model()
    
    # Header
    st.title("🔍 FAISS Vector Search Educational Demo")
    st.markdown("Learn how vector similarity search works by adding documents and searching through them!")
    
    # Display educational panel
    display_educational_panel()
    
    # Main content in columns
    col1, col2 = st.columns([1, 1])
    
    # Section 1: Add Documents
    with col1:
        st.header("📄 Document Management")
        
        # Create tabs for different input methods
        tab1, tab2 = st.tabs(["✍️ Manual Input", "📁 File Upload"])
        
        # Tab 1: Manual text input
        with tab1:
            new_doc = st.text_area(
                "✍️ Enter a new document:",
                placeholder="Type your document here... (e.g., 'Machine learning helps computers learn from data')",
                help="Add text documents that will be converted to embeddings and stored in FAISS",
                height=100,
                key="manual_input"
            )
            
            if st.button("➕ Add Text Document", type="primary", use_container_width=True, key="add_manual"):
                if new_doc.strip():
                    with st.spinner("🔄 Generating embedding..."):
                        # Generate embedding
                        embedding = embed_text(new_doc.strip(), st.session_state.model)
                        
                        # Add to session state
                        st.session_state.documents.append(new_doc.strip())
                        st.session_state.embeddings.append(embedding)
                        
                        # Recreate FAISS index
                        st.session_state.faiss_index = create_faiss_index(st.session_state.embeddings)
                        
                        st.success(f"✅ Document added! Total: {len(st.session_state.documents)} documents")
                        st.rerun()
                else:
                    st.error("⚠️ Please enter some text!")
        
        # Tab 2: File upload
        with tab2:
            # Initialize uploaded_files to None
            uploaded_files = None
            
            # Determine available file types based on installed libraries
            available_types = []
            help_text_parts = ["Supported formats: "]
            
            if PDF_AVAILABLE:
                available_types.append('pdf')
                help_text_parts.append("PDF")
            
            if DOCX_AVAILABLE:
                available_types.extend(['docx', 'doc'])
                help_text_parts.append("Word documents (.docx, .doc)")
            
            if PPTX_AVAILABLE:
                available_types.extend(['pptx', 'ppt'])
                help_text_parts.append("PowerPoint (.pptx, .ppt)")
            
            if not available_types:
                st.warning("⚠️ No file processing libraries available. Only manual text input is supported.")
                st.info("To enable file uploads, install: pypdf2, python-docx, python-pptx")
            else:
                help_text = ", ".join(help_text_parts[1:])
                
                uploaded_files = st.file_uploader(
                    "📁 Upload documents:",
                    type=available_types,
                    accept_multiple_files=True,
                    help=help_text_parts[0] + help_text
                )
            
            # SAFETY: Ensure uploaded_files is always defined
            if 'uploaded_files' not in locals():
                uploaded_files = None
                
            if uploaded_files:
                for uploaded_file in uploaded_files:
                    if st.button(f"➕ Process {uploaded_file.name}", key=f"process_{uploaded_file.name}"):
                        with st.spinner(f"🔄 Processing {uploaded_file.name}..."):
                            # Extract text from file
                            extracted_text, file_info = process_uploaded_file(uploaded_file)
                            
                            if extracted_text and not extracted_text.startswith("Error") and not extracted_text.startswith("Unsupported"):
                                # Generate embedding
                                embedding = embed_text(extracted_text, st.session_state.model)
                                
                                # Create document entry with file info
                                doc_entry = f"[FILE: {uploaded_file.name}]\n{extracted_text[:500]}{'...' if len(extracted_text) > 500 else ''}"
                                
                                # Add to session state
                                st.session_state.documents.append(doc_entry)
                                st.session_state.embeddings.append(embedding)
                                
                                # Recreate FAISS index
                                st.session_state.faiss_index = create_faiss_index(st.session_state.embeddings)
                                
                                st.success(f"✅ File processed! {file_info}")
                                st.success(f"Total documents: {len(st.session_state.documents)}")
                                st.rerun()
                            else:
                                st.error(f"⚠️ {extracted_text}")
        
        # Show current documents
        if st.session_state.documents:
            st.subheader(f"📚 Stored Documents ({len(st.session_state.documents)})")
            
            for i, doc in enumerate(st.session_state.documents):
                # Check if it's a file-based document
                is_file_doc = doc.startswith('[FILE:')
                
                if is_file_doc:
                    # Extract filename from document
                    filename = doc.split(']')[0].replace('[FILE: ', '')
                    preview_text = doc.split('\n', 1)[1] if '\n' in doc else doc
                    expander_title = f"📁 {filename}"
                else:
                    preview_text = doc
                    expander_title = f"📄 Document {i+1}: {doc[:30]}..."
                
                with st.expander(expander_title):
                    if is_file_doc:
                        st.write(f"**File:** {filename}")
                        st.write(f"**Content preview:** {preview_text[:200]}..." if len(preview_text) > 200 else preview_text)
                    else:
                        st.write(f"**Text:** {doc}")
                    
                    with st.expander("🔢 View Embedding Preview"):
                        st.code(f"Vector: [{', '.join([f'{x:.3f}' for x in st.session_state.embeddings[i][:8]])}...] (384 dimensions)", language="python")
                    
                    if st.button(f"🗑️ Remove Document {i+1}", key=f"remove_{i}"):
                        st.session_state.documents.pop(i)
                        st.session_state.embeddings.pop(i)
                        st.session_state.faiss_index = create_faiss_index(st.session_state.embeddings)
                        st.session_state.search_results = None
                        st.rerun()
    
    # Section 2: Search
    with col2:
        st.header("🔍 Search Documents")
        
        query = st.text_input(
            "Enter search query:",
            placeholder="What are you looking for?",
            help="Search through your documents using vector similarity"
        )
        
        k = st.slider("Number of results:", 1, min(10, len(st.session_state.documents)) if st.session_state.documents else 1, 3)
        
        if st.button("Search", type="primary") and query.strip():
            if not st.session_state.documents:
                st.error("Please add some documents first!")
            else:
                with st.spinner("Searching..."):
                    # Generate query embedding
                    query_embedding = embed_text(query.strip(), st.session_state.model)
                    
                    # Search FAISS
                    scores, indices = search_faiss(query_embedding, st.session_state.faiss_index, k)
                    
                    # Store results
                    st.session_state.search_results = {
                        'query': query.strip(),
                        'query_embedding': query_embedding,
                        'scores': scores,
                        'indices': indices
                    }
                    
                    st.success("Search completed!")
        
        # Display search results
        if st.session_state.search_results:
            results = st.session_state.search_results
            
            st.subheader("🎯 Search Results")
            
            # Query embedding preview
            with st.expander("Query Embedding"):
                st.write(f"**Query:** {results['query']}")
                st.write(f"**Embedding preview:** {results['query_embedding'][:5]}... (384 dimensions)")
            
            # Results
            st.write("**Ranked Results:**")
            for rank, (score, idx) in enumerate(zip(results['scores'], results['indices'])):
                if idx < len(st.session_state.documents):  # Safety check
                    with st.container():
                        st.write(f"**#{rank+1} - Similarity: {score:.4f}**")
                        st.write(f"📄 {st.session_state.documents[idx]}")
                        st.write("---")
    
    # Section 3: Visualization
    if st.session_state.documents:
        st.header("📊 Vector Space Visualization")
        
        # Controls
        col_viz1, col_viz2 = st.columns([1, 3])
        
        with col_viz1:
            reduction_method = st.selectbox(
                "Dimension reduction method:",
                ['pca', 'tsne'],
                help="PCA is faster, t-SNE may show better clusters"
            )
            
            st.metric("Total Documents", len(st.session_state.documents))
            
            if st.session_state.faiss_index:
                st.metric("Vector Dimension", st.session_state.embeddings[0].shape[0])
        
        with col_viz2:
            # Generate 2D representation
            with st.spinner(f"Reducing dimensions using {reduction_method.upper()}..."):
                embeddings_2d = reduce_dimensions(st.session_state.embeddings, reduction_method)
                
                # Prepare visualization data
                query_2d = None
                similar_indices = None
                similarity_scores = None
                
                if st.session_state.search_results:
                    # Add query to embeddings for reduction
                    all_embeddings = st.session_state.embeddings + [st.session_state.search_results['query_embedding']]
                    all_2d = reduce_dimensions(all_embeddings, reduction_method)
                    
                    embeddings_2d = all_2d[:-1]  # Document embeddings
                    query_2d = all_2d[-1]       # Query embedding
                    similar_indices = list(st.session_state.search_results['indices'])
                    similarity_scores = st.session_state.search_results['scores']
                
                # Create interactive plot
                fig = plot_embeddings_interactive(
                    embeddings_2d, 
                    st.session_state.documents,
                    query_2d,
                    similar_indices,
                    similarity_scores
                )
                
                st.plotly_chart(fig, use_container_width=True)
        
        # Explanation
        st.info("""
        **Understanding the Visualization:**
        - Each point represents a document in the high-dimensional embedding space
        - Closer points are more similar in meaning
        - When you search, FAISS finds the closest points to your query (green star)
        - The red points are the most similar documents found
        """)
    
    else:
        st.header("📊 Vector Space Visualization")
        st.info("Add some documents to see the vector space visualization!")
    
    # Footer with statistics
    if st.session_state.documents:
        st.markdown("---")
        col_stats1, col_stats2, col_stats3 = st.columns(3)
        
        with col_stats1:
            st.metric("Documents in Index", len(st.session_state.documents))
        
        with col_stats2:
            if st.session_state.search_results:
                avg_score = np.mean(st.session_state.search_results['scores'])
                st.metric("Avg Similarity Score", f"{avg_score:.4f}")
            else:
                st.metric("Searches Performed", "0")
        
        with col_stats3:
            total_chars = sum(len(doc) for doc in st.session_state.documents)
            st.metric("Total Characters", f"{total_chars:,}")

if __name__ == "__main__":
    main()